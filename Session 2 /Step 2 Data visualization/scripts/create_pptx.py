#!/usr/bin/env python3
"""
Create a WHU-styled PPTX from a JSON spec, using assets/whu-overall-template.pptx.
This script focuses on reliable placeholder filling, not on rebuilding layouts.

Usage:
  python scripts/create_pptx.py --spec assets/example_pptx_spec.json --out out.pptx
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from pptx import Presentation

TEMPLATE_REL = Path(__file__).resolve().parent.parent / "assets" / "whu-overall-template.pptx"

def _find_layout(prs: Presentation, layout_name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.strip() == layout_name.strip():
                return layout
    raise ValueError(f"Unknown layout: {layout_name!r}. Check references/LAYOUTS.json for valid names.")

def _ph(slide, idx: int):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == idx:
                return ph
        except Exception:
            continue
    raise KeyError(f"Placeholder idx={idx} not found on slide (layout={slide.slide_layout.name}).")

def _set_text(ph, text: str | None):
    if text is None:
        return
    # Keep placeholder formatting: set text via text_frame without redefining fonts
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text

def build(spec_path: Path, out_path: Path):
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    prs = Presentation(str(TEMPLATE_REL))

    # start from an empty deck: remove all existing slides by creating a new Presentation from template theme
    # python-pptx cannot "delete" slides cleanly, so we keep the template slides but add new ones at the end.
    # Consumers can delete the template example slides manually if needed.
    for s in spec.get("slides", []):
        layout = _find_layout(prs, s["layout"])
        slide = prs.slides.add_slide(layout)

        # Common placeholder conventions in this WHU template:
        # 11 headline, 12 topic, 13 subheadline, 14 body, 4 slide number, 10 date, 13/14/15/16 cover lines
        # Fill only keys present in the slide spec.
        if "headline" in s:
            try: _set_text(_ph(slide, 11), s["headline"])
            except Exception: pass
        if "topic" in s:
            try: _set_text(_ph(slide, 12), s["topic"])
            except Exception: pass
        if "subheadline" in s:
            try: _set_text(_ph(slide, 13), s["subheadline"])
            except Exception: pass
        if "body" in s:
            try: _set_text(_ph(slide, 14), s["body"])
            except Exception: pass

        # Bullets: write into body placeholder if provided and body is empty
        bullets = s.get("bullets") or []
        if bullets:
            try:
                ph = _ph(slide, 14)
                tf = ph.text_frame
                tf.clear()
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(b)
                    p.level = 0
            except Exception:
                pass

        # Cover slide support (WHU_4-Part Headline and similar)
        if "cover_part1" in s:
            for idx, key in [(14,"cover_part1"), (15,"cover_part2"), (13,"author_line"), (10,"date_line")]:
                if key in s:
                    try: _set_text(_ph(slide, idx), s[key])
                    except Exception: pass

    prs.save(str(out_path))

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True, type=Path, help="Path to JSON spec")
    ap.add_argument("--out", required=True, type=Path, help="Output PPTX path")
    args = ap.parse_args()
    build(args.spec, args.out)

if __name__ == "__main__":
    main()
